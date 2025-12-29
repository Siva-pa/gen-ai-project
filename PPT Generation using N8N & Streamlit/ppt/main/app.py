import streamlit as st
import requests
import json
import re
from pptx import Presentation

# --------------------------------------------------
# PAGE CONFIG
# --------------------------------------------------
st.set_page_config(page_title="AI PPT Generator", layout="centered")

st.title("📄 AI PowerPoint Generator")
st.subheader("Generate PPT slides using AI 🤖")

# --------------------------------------------------
# USER INPUT
# --------------------------------------------------
user_input = st.text_area(
    "Enter your topic or content",
    placeholder="Example: Power BI"
)

num_slides = st.slider(
    "Number of slides",
    min_value=1,
    max_value=15,
    value=3
)

# --------------------------------------------------
# HELPER FUNCTION: CLEAN AI RESPONSE
# --------------------------------------------------
def clean_ai_json(text: str) -> str:
    """
    Removes markdown ```json wrappers and extra spaces.
    """
    text = text.strip()

    # Remove markdown fences
    text = re.sub(r"^```json\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^```\s*", "", text)
    text = re.sub(r"\s*```$", "", text)

    return text.strip()

# --------------------------------------------------
# GENERATE PPT
# --------------------------------------------------
if st.button("Generate PPT"):
    if not user_input.strip():
        st.warning("⚠️ Please enter a topic.")
        st.stop()

    with st.spinner("Generating slides using AI..."):
        response = requests.post(
            url="https://anna-sovice1.app.n8n.cloud/webhook/bfeef820-0a2e-47c7-a74c-f41cde078ec9",
            json={
                "user_input": user_input.strip(),
                "num_slides": num_slides
            },
            timeout=120
        )

    # --------------------------------------------------
    # CHECK EMPTY RESPONSE
    # --------------------------------------------------
    if not response.text.strip():
        st.error("❌ Empty response from AI service.")
        st.stop()

    # --------------------------------------------------
    # CLEAN & PARSE JSON
    # --------------------------------------------------
    cleaned_text = clean_ai_json(response.text)

    try:
        data = json.loads(cleaned_text)
    except json.JSONDecodeError:
        st.error("❌ Invalid JSON received from AI service.")
        st.write("Raw response:", response.text)
        st.write("Cleaned response:", cleaned_text)
        st.stop()

    # --------------------------------------------------
    # CREATE POWERPOINT
    # --------------------------------------------------
    prs = Presentation()

    # ---------- TITLE SLIDE ----------
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = user_input.strip().title()

    # ---------- CONTENT SLIDES ----------
    slides = data.get("slides", [])[:num_slides]

    for idx, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # ---- SAFE TITLE ----
        slide_title = slide_data.get("title", "").strip()
        if not slide_title:
            slide_title = f"{user_input.title()} – Slide {idx}"

        slide.shapes.title.text = slide_title

        # ---- SAFE CONTENT ----
        content = slide_data.get("content", [])
        valid_bullets = [c.strip() for c in content if c.strip()]

        if not valid_bullets:
            valid_bullets = [
                f"Overview of {user_input}",
                f"Key features of {user_input}",
                f"Why {user_input} is important"
            ]

        slide.placeholders[1].text = "\n".join(valid_bullets)

    # --------------------------------------------------
    # SAVE & DOWNLOAD
    # --------------------------------------------------
    file_name = "AI_Presentation.pptx"
    prs.save(file_name)

    with open(file_name, "rb") as f:
        st.success("✅ PPT generated successfully!")
        st.download_button(
            label="⬇️ Download PPT",
            data=f,
            file_name=file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )